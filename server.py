"""
server.py — Single-process FastAPI WebSocket server with live ASR streaming.

Architecture (D002):
    One FastAPI app serves:
      - ``GET /`` → ``static/index.html`` (browser frontend)
      - ``WS /ws`` → live transcription pipeline via ``Transcriber`` (D004, MEM005)

    WebSocket lifecycle maps 1:1 to ASR streaming lifecycle:
      connect → ``start_streaming()``
      binary frame → ``stream_chunk()``
      disconnect → ``finish_streaming()``

Single-session constraint:
    Only one WebSocket may be active at a time because the Transcriber holds a
    single Qwen3ASRModel.LLM() instance.  If a second WebSocket connects while
    one is active, the server logs a warning but does **not** reject the
    connection — however, the first connection's streaming state will be
    corrupted since both share the same Transcriber.  This is documented in
    the log and is acceptable for a single-user local tool.
"""

from __future__ import annotations

import io
import json
import logging
import os
import sys
import time
import traceback
from contextlib import asynccontextmanager
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

import numpy as np
from fastapi import FastAPI, File, UploadFile, WebSocket, WebSocketDisconnect
from fastapi.responses import JSONResponse, Response
from fastapi.staticfiles import StaticFiles
from starlette.responses import FileResponse

try:
    from transcriber import Transcriber
    _gpu_available = True
except ImportError:
    Transcriber = None  # type: ignore[assignment]
    _gpu_available = False

from transcriber_cpu import TranscriberCPU

# ── Structured logger (same pattern as transcriber.py) ──────────────────────
logger = logging.getLogger("server")
logger.setLevel(logging.INFO)
_handler = logging.StreamHandler(sys.stderr)
_handler.setFormatter(
    logging.Formatter(
        "%(asctime)s [%(levelname)s] %(name)s: %(message)s",
        datefmt="%Y-%m-%dT%H:%M:%S",
    )
)
logger.addHandler(_handler)


# ── Module-level state ──────────────────────────────────────────────────────
_transcriber = None
"""Module-level Transcriber instance.  If model load fails, this stays ``None``
and ``/ws`` returns 503."""

_model_available: bool = False
"""Set to ``True`` after successful model load in the lifespan handler."""

_cpu_mode: bool = os.environ.get("TRANSCRIBER_CPU", "") == "1"
"""Whether the server is running in CPU mode (using TranscriberCPU)."""

_cpu_model_dir: str = os.environ.get("TRANSCRIBER_CPU_MODEL_DIR", "qwen3-asr-0.6b")
"""Model directory for CPU mode."""

_cpu_binary_path: str = os.environ.get("TRANSCRIBER_CPU_BINARY", "./qwen-asr/qwen_asr")
"""Path to the qwen_asr binary for CPU mode."""

_active_streaming: bool = False
"""Tracks whether a WebSocket streaming session is currently active.
Only one session at a time (single LLM instance)."""

_session_active: bool = False
"""Whether a recording session is currently in progress.  Set ``True`` on
first audio chunk, ``False`` on stop or disconnect."""

_session_paused: bool = False
"""Whether the session is paused.  While ``True``, audio frames are silently
skipped to prevent ASR state corruption."""

_session_transcript: List[Dict[str, Any]] = []
"""Accumulated transcript segments.  Each segment is a dict with at least
``text`` (str) and ``time`` (float, wall-clock seconds since session start)."""

_session_context: str = ""
"""Domain-context text extracted from uploaded documents (R003)."""

_session_start_time: Optional[float] = None
"""Wall-clock timestamp (``time.time()``) when the current session began.
``None`` when no session is active."""


# ── Lifespan ────────────────────────────────────────────────────────────────


@asynccontextmanager
async def lifespan(app: FastAPI):
    """Load the Transcriber model at startup, clean up on shutdown."""
    global _transcriber, _model_available

    if _cpu_mode:
        logger.info("Server starting in CPU mode — loading TranscriberCPU…")
        try:
            _transcriber = TranscriberCPU(
                model_dir=_cpu_model_dir,
                binary_path=_cpu_binary_path,
                language="English",
            )
            _model_available = True
            logger.info("CPU ASR engine ready. Server ready.")
        except Exception as exc:
            logger.error("CPU model load failed: %s", exc)
            logger.error(
                "Server will start, but /ws will return 503 until the model "
                "is available."
            )
            _model_available = False
    else:
        logger.info("Server starting — loading GPU ASR model…")
        try:
            if Transcriber is None:
                raise ImportError("GPU backend (transcriber) not available")
            _transcriber = Transcriber()
            _model_available = True
            logger.info("ASR model loaded successfully. Server ready.")
        except (ImportError, RuntimeError) as exc:
            logger.error("Model load failed: %s", exc)
            logger.error(
                "Server will start, but /ws will return 503 until the model "
                "is available."
            )
            _model_available = False
        except Exception:
            logger.error(
                "Unexpected error during model load:\n%s", traceback.format_exc()
            )
            _model_available = False

    yield  # ── Server runs here ──

    logger.info("Server shutting down.")


# ── FastAPI app ─────────────────────────────────────────────────────────────

app = FastAPI(title="Transcriber", lifespan=lifespan)


@app.get("/")
async def index() -> FileResponse:
    """Serve the browser frontend."""
    return FileResponse("static/index.html")


# Mount static file serving for the browser frontend (T02)
app.mount("/static", StaticFiles(directory="static"), name="static")


# ── WebSocket endpoint ──────────────────────────────────────────────────────


@app.websocket("/ws")
async def websocket_transcribe(ws: WebSocket) -> None:
    """Live transcription WebSocket endpoint.

    Protocol:
        - Client sends binary audio frames (float32, 16 kHz mono PCM).
        - Server returns JSON text frames: ``{"text": "...", "language": "English"}``.
        - On error, server returns: ``{"error": "…"}``.

    Lifecycle:
        connect → ``start_streaming()`` → binary frames → ``stream_chunk()``
        → disconnect → ``finish_streaming()``.
    """
    global _active_streaming, _transcriber, _model_available
    global _session_active, _session_paused, _session_start_time

    client_host = ws.client.host if ws.client else "unknown"
    client_port = ws.client.port if ws.client else 0

    # ── Check model availability ────────────────────────────────────────
    if not _model_available or _transcriber is None:
        logger.warning(
            "/ws connection from %s:%d rejected — model not available.",
            client_host,
            client_port,
        )
        await ws.accept()
        await ws.send_json({"error": "Model not available"})
        await ws.close(code=1011)
        return

    # ── Accept the WebSocket ─────────────────────────────────────────────
    await ws.accept()
    logger.info(
        "WebSocket connected: %s:%d — streaming session started.",
        client_host,
        client_port,
    )

    # ── Single-session constraint ────────────────────────────────────────
    if _active_streaming:
        logger.warning(
            "Second WebSocket connected while a streaming session is already "
            "active.  Only one session is supported — the first connection's "
            "state may be corrupted."
        )
    _active_streaming = True

    state: Any = None
    chunk_count = 0
    total_text_chars = 0

    try:
        # ── Start streaming ──────────────────────────────────────────────
        state = _transcriber.start_streaming(chunk_size_sec=2.0)
        # vLLM's init_streaming_state returns an ASRStreamingState object
        # with .text and .language attributes.
        initial_text = state.text
        if initial_text:
            await ws.send_json({"text": initial_text, "language": _transcriber.language})
            logger.info("Initial text sent (%d chars).", len(initial_text))

        # ── Receive loop ─────────────────────────────────────────────────
        while True:
            frame = await ws.receive()

            # ── Text frame: control protocol ────────────────────────────
            if "text" in frame:
                try:
                    msg = json.loads(frame["text"])
                except json.JSONDecodeError:
                    logger.warning(
                        "Invalid JSON in control message from %s:%d.",
                        client_host,
                        client_port,
                    )
                    await ws.send_json({"error": "Invalid JSON in control message"})
                    continue

                action = msg.get("action", "")

                # ── "pause": finalise current segment, save transcript ─
                if action == "pause":
                    if state is not None and not _session_paused:
                        result = _transcriber.finish_streaming(state)
                        _session_transcript.append(
                            {"text": result.text, "time": time.time()}
                        )
                        state = None
                    _session_paused = True
                    segment_text = (
                        _session_transcript[-1]["text"] if _session_transcript else ""
                    )
                    await ws.send_json(
                        {"status": "paused", "text": segment_text}
                    )
                    logger.info(
                        "Pause: session paused — %d segments accumulated.",
                        len(_session_transcript),
                    )

                # ── "resume": start a new ASR streaming segment ─────────
                elif action == "resume":
                    state = _transcriber.start_streaming(chunk_size_sec=2.0)
                    initial_text = state.text
                    if initial_text:
                        await ws.send_json(
                            {"text": initial_text, "language": _transcriber.language}
                        )
                    _session_paused = False
                    await ws.send_json({"status": "resumed"})
                    logger.info("Resume: session resumed.")

                # ── "stop": finalise and end session ─────────────────────
                elif action == "stop":
                    final_text = ""
                    if state is not None:
                        result = _transcriber.finish_streaming(state)
                        final_text = result.text
                        _session_transcript.append(
                            {"text": result.text, "time": time.time()}
                        )
                        state = None
                        await ws.send_json(
                            {"text": result.text, "language": _transcriber.language}
                        )
                    _session_active = False
                    _session_paused = False
                    await ws.send_json({
                        "status": "stopped",
                        "text": final_text,
                        "chunks": chunk_count,
                        "segments": len(_session_transcript),
                    })
                    logger.info(
                        "Stop: session stopped — %d segments accumulated.",
                        len(_session_transcript),
                    )
                    break

                # ── Unknown action ───────────────────────────────────────
                else:
                    await ws.send_json(
                        {"error": f"Unknown action: {action}"}
                    )
                    logger.warning(
                        "Unknown action '%s' from %s:%d.",
                        action,
                        client_host,
                        client_port,
                    )

                continue

            # ── Binary frame: audio processing ──────────────────────────
            if "bytes" not in frame:
                logger.warning(
                    "Unexpected WebSocket frame type %s from %s:%d — ignoring.",
                    set(frame.keys()) - {"type"},
                    client_host,
                    client_port,
                )
                continue

            raw_bytes = frame["bytes"]

            # Pause guard: silently skip audio frames while paused
            if _session_paused:
                logger.debug("Audio frame skipped — session paused.")
                continue

            # First audio chunk — mark session active
            if not _session_active:
                _session_active = True
                _session_start_time = time.time()
                logger.info("Recording session active — first audio chunk received.")

            # Guard: skip zero-length frames
            if len(raw_bytes) == 0:
                logger.debug(
                    "Skipping zero-length binary frame (chunk %d).",
                    chunk_count + 1,
                )
                continue

            chunk_count += 1

            # Convert bytes to float32 array
            try:
                audio_chunk = np.frombuffer(raw_bytes, dtype=np.float32)
            except ValueError:
                logger.warning(
                    "Chunk %d: byte payload size %d is not a multiple of "
                    "float32 size — np.frombuffer will handle remainder. "
                    "Skipping frame.",
                    chunk_count,
                    len(raw_bytes),
                )
                await ws.send_json({"error": "Invalid audio frame: byte size "
                                              "not aligned to float32"})
                continue

            # Guard: skip empty or tiny audio
            if len(audio_chunk) == 0:
                logger.debug(
                    "Skipping empty audio chunk (chunk %d).", chunk_count
                )
                continue

            # Process the chunk through the Transcriber
            try:
                text = _transcriber.stream_chunk(audio_chunk, state)
            except Exception:
                logger.error(
                    "Error processing audio chunk %d:\n%s",
                    chunk_count,
                    traceback.format_exc(),
                )
                await ws.send_json({"error": "Transcription processing error"})
                continue

            # Send partial text back to the browser
            if text:
                total_text_chars = len(text)
                await ws.send_json(
                    {"text": text, "language": _transcriber.language}
                )
                logger.debug(
                    "Chunk %d: sent %d chars.", chunk_count, len(text)
                )

    except WebSocketDisconnect:
        logger.info(
            "WebSocket disconnected by client: %s:%d. "
            "Flushing final tokens…",
            client_host,
            client_port,
        )
    except Exception:
        logger.error(
            "Unexpected error during streaming session %s:%d:\n%s",
            client_host,
            client_port,
            traceback.format_exc(),
        )
    finally:
        # ── Finish streaming (flush final tokens) ────────────────────────
        if state is not None and _transcriber is not None:
            try:
                result = _transcriber.finish_streaming(state)
                total_text_chars = len(result.text)
                # Save final segment if session was active
                if _session_active:
                    _session_transcript.append(
                        {"text": result.text, "time": time.time()}
                    )
                # Send final text if it differs from last partial
                try:
                    await ws.send_json(
                        {"text": result.text, "language": result.language}
                    )
                except Exception:
                    # Client likely already disconnected — final text is logged.
                    pass
            except Exception:
                logger.error(
                    "Error during finish_streaming:\n%s",
                    traceback.format_exc(),
                )

        _active_streaming = False
        _session_active = False
        _session_paused = False
        logger.info(
            "Streaming session ended: %s:%d — %d chunks processed, "
            "final text: %d chars.",
            client_host,
            client_port,
            chunk_count,
            total_text_chars,
        )


# ── Health check ────────────────────────────────────────────────────────────


@app.get("/health")
async def health() -> Dict[str, Any]:
    """Health endpoint — returns model availability and active status."""
    return {
        "status": "ok",
        "model_available": _model_available,
        "active_streaming": _active_streaming,
        "session_active": _session_active,
        "session_paused": _session_paused,
        "cpu_mode": _cpu_mode,
    }


# ── Document upload (R003) ──────────────────────────────────────────────────


ALLOWED_EXTENSIONS = frozenset({".txt", ".docx", ".pptx"})


@app.post("/upload")
async def upload_document(file: UploadFile = File(...)) -> JSONResponse:
    """Upload a document for domain-context hinting (R003).

    Supported formats: ``.txt``, ``.docx``, ``.pptx``.

    The extracted text is stored in ``_session_context`` and, if the
    Transcriber is available, pushed via ``update_context()`` so that
    subsequent transcription passes benefit from domain-term hints.
    """
    global _session_context, _transcriber

    # ── Validate file extension ──────────────────────────────────────────
    if not file.filename:
        logger.warning("Upload rejected: no filename provided.")
        return JSONResponse(
            status_code=400,
            content={"error": "No filename provided."},
        )

    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        logger.warning(
            "Upload rejected: unsupported extension '%s' (file=%s).",
            ext,
            file.filename,
        )
        return JSONResponse(
            status_code=400,
            content={
                "error": (
                    f"Unsupported file type: {ext}. "
                    "Allowed: .txt, .docx, .pptx"
                )
            },
        )

    # ── Read file content ────────────────────────────────────────────────
    try:
        content = await file.read()
    except Exception as exc:
        logger.error("Failed to read uploaded file '%s': %s", file.filename, exc)
        return JSONResponse(
            status_code=400,
            content={"error": f"Failed to read file: {exc}"},
        )

    if not content:
        logger.warning("Upload rejected: empty file '%s'.", file.filename)
        return JSONResponse(
            status_code=400,
            content={"error": "Uploaded file is empty."},
        )

    # ── Extract text by type ─────────────────────────────────────────────
    try:
        text = _extract_text(content, ext, file.filename)
    except Exception as exc:
        logger.error(
            "Failed to parse file '%s': %s",
            file.filename,
            exc,
            exc_info=True,
        )
        return JSONResponse(
            status_code=400,
            content={"error": f"Failed to parse file: {exc}"},
        )

    if not text.strip():
        logger.warning(
            "Uploaded file '%s' yielded no extractable text.", file.filename
        )

    # ── Store context ────────────────────────────────────────────────────
    _session_context = text
    if _transcriber is not None:
        _transcriber.update_context(text)

    logger.info(
        "Document uploaded: %s (%d chars, %s). Context stored.",
        file.filename,
        len(text),
        ext,
    )

    return JSONResponse(
        status_code=200,
        content={"status": "ok", "chars": len(text), "filename": file.filename},
    )


def _extract_text(content: bytes, ext: str, filename: str) -> str:
    """Extract plain text from a document by extension.

    Raises:
        ImportError: If the required library is not installed.
        ValueError: If the file content cannot be parsed.
    """
    if ext == ".txt":
        return content.decode("utf-8")

    elif ext == ".docx":
        try:
            from docx import Document  # type: ignore[import-untyped]
        except ImportError:
            raise ImportError(
                "python-docx is not installed. Install with: "
                "pip install python-docx>=1.0,<2.0"
            )
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)

    elif ext == ".pptx":
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
        except ImportError:
            raise ImportError(
                "python-pptx is not installed. Install with: "
                "pip install python-pptx>=0.6,<1.0"
            )
        prs = Presentation(io.BytesIO(content))
        texts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    slide_text = shape.text_frame.text.strip()
                    if slide_text:
                        texts.append(slide_text)
        return "\n".join(texts)

    # Should not reach here — caller validates extension.
    raise ValueError(f"Unsupported file extension: {ext}")


# ── Transcript export (R005) ────────────────────────────────────────────────


@app.get("/export")
async def export_transcript() -> Response:
    """Export the accumulated transcript as timestamped markdown (R005).

    Returns a downloadable ``.md`` file with an ISO-8601 date header,
    language and context metadata, and timestamped segments in
    ``[HH:MM:SS] text`` format.  Works with an empty transcript
    (header-only output).
    """
    global _session_transcript, _transcriber

    now = datetime.now(timezone.utc)

    # ── Build markdown ───────────────────────────────────────────────────
    lines: List[str] = []
    lines.append("# Meeting Transcript")
    lines.append("")
    lines.append(f"**Date:** {now.strftime('%Y-%m-%dT%H:%M:%SZ')}")
    lines.append("")

    if _transcriber is not None:
        lines.append(f"**Language:** {_transcriber.language}")
        if _transcriber.context:
            lines.append(f"**Context:** {_transcriber.context}")
    else:
        lines.append("**Language:** English")

    lines.append("")

    if _session_transcript:
        lines.append("---")
        lines.append("")

        base_time: Optional[float] = _session_start_time

        for segment in _session_transcript:
            seg_text = segment.get("text", "")
            seg_time = segment.get("time", 0.0)

            # If we have a base time, use it; otherwise treat seg_time as
            # already relative.
            if base_time is not None:
                elapsed = seg_time - base_time
            else:
                elapsed = seg_time

            # Format as HH:MM:SS
            hours = int(elapsed // 3600)
            minutes = int((elapsed % 3600) // 60)
            seconds = int(elapsed % 60)
            timestamp = f"{hours:02d}:{minutes:02d}:{seconds:02d}"

            lines.append(f"[{timestamp}] {seg_text}")
            lines.append("")

    # ── Build filename ───────────────────────────────────────────────────
    date_str = now.strftime("%Y-%m-%d")
    filename = f"transcript-{date_str}.md"

    md = "\n".join(lines)
    logger.info(
        "Export: %d segments, %d chars → %s.",
        len(_session_transcript),
        len(md),
        filename,
    )

    return Response(
        content=md,
        media_type="text/markdown; charset=utf-8",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


# ── Entrypoint ──────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import argparse
    import uvicorn

    parser = argparse.ArgumentParser(description="Transcriber Server")
    parser.add_argument(
        "--cpu",
        action="store_true",
        help="Run in CPU mode using antirez/qwen-asr binary (no GPU required).",
    )
    parser.add_argument(
        "--cpu-model-dir",
        default="qwen3-asr-0.6b",
        help="Model directory for CPU mode (default: qwen3-asr-0.6b).",
    )
    parser.add_argument(
        "--cpu-binary-path",
        default="./qwen-asr/qwen_asr",
        help="Path to qwen_asr binary (default: ./qwen-asr/qwen_asr).",
    )
    parser.add_argument(
        "--host",
        default="0.0.0.0",
        help="Bind address (default: 0.0.0.0).",
    )
    parser.add_argument(
        "--port",
        type=int,
        default=5000,
        help="Bind port (default: 5000).",
    )
    args = parser.parse_args()

    if args.cpu:
        os.environ["TRANSCRIBER_CPU"] = "1"
        os.environ["TRANSCRIBER_CPU_MODEL_DIR"] = args.cpu_model_dir
        os.environ["TRANSCRIBER_CPU_BINARY"] = args.cpu_binary_path
        print(f"CPU mode: model={args.cpu_model_dir}, binary={args.cpu_binary_path}")

    uvicorn.run(
        "server:app",
        host=args.host,
        port=args.port,
        log_level="info",
        reload=False,
    )
